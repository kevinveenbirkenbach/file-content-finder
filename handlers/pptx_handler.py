# pptx_handler.py
from pptx import Presentation
from .base_handler import BaseHandler
from models import FileResult

class PPTXHandler(BaseHandler):
    def search(self):
        find_cmd = ['find', self.search_path, '-type', 'f', '-iname', self.file_type, '-print0']
        return self.process_files_in_parallel(find_cmd, self.process_pptx)

    def process_pptx(self, file_path):
        results = []
        try:
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text
            for search_string in self.search_strings:
                if search_string in text:
                    results.append(FileResult(file_path, self.file_type, text))
        except Exception as e:
            self.error_handler(str(e), file_path)
        return results
